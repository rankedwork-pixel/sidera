"""Build the Project Sidera pitch deck as a .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand colours ──────────────────────────────────────────────
BG = RGBColor(0x0A, 0x0A, 0x0F)
BG_CARD = RGBColor(0x14, 0x14, 0x1E)
TEXT = RGBColor(0xE8, 0xE8, 0xF0)
TEXT_DIM = RGBColor(0x88, 0x88, 0xA0)
ACCENT = RGBColor(0x08, 0x75, 0x6C)
ACCENT_LIGHT = RGBColor(0x0A, 0xA8, 0x9B)
RED = RGBColor(0xE7, 0x4C, 0x5E)
ORANGE = RGBColor(0xF0, 0x97, 0x3B)
BLUE = RGBColor(0x4A, 0x9E, 0xFF)
GREEN = RGBColor(0x34, 0xD3, 0x99)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Paths
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
REPO_DIR = os.path.dirname(SCRIPT_DIR)
LOGO_PATH = os.path.join(REPO_DIR, "SLACK IMAGE_square.png")


# ── Helpers ────────────────────────────────────────────────────
def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_para(text_frame, text, font_size=18, color=TEXT, bold=False,
             alignment=PP_ALIGN.LEFT, space_before=Pt(6), font_name="Calibri"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD,
                     border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def card_with_text(slide, left, top, width, height, title, body,
                   title_color=ACCENT_LIGHT, border_color=None):
    card = add_rounded_rect(slide, left, top, width, height,
                            border_color=border_color)
    # Title
    add_text_box(slide, left + Inches(0.25), top + Inches(0.2),
                 width - Inches(0.5), Inches(0.5),
                 title, font_size=20, color=title_color, bold=True)
    # Body
    add_text_box(slide, left + Inches(0.25), top + Inches(0.7),
                 width - Inches(0.5), height - Inches(0.9),
                 body, font_size=14, color=TEXT_DIM)
    return card


def add_slide_number(slide, num, total):
    add_text_box(slide, Inches(0.5), SLIDE_H - Inches(0.5),
                 Inches(1.5), Inches(0.3),
                 f"{num} / {total}", font_size=10, color=TEXT_DIM,
                 font_name="Consolas")


def add_watermark(slide):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(
            LOGO_PATH, SLIDE_W - Inches(1.1), Inches(0.25),
            width=Inches(0.4), height=Inches(0.4)
        )


TOTAL_SLIDES = 27


# ── SLIDE 1: Title ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)
if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(
        LOGO_PATH,
        SLIDE_W / 2 - Inches(1.2), Inches(1.0),
        width=Inches(2.4), height=Inches(2.4)
    )
add_text_box(slide, Inches(0), Inches(3.6), SLIDE_W, Inches(1),
             "Project Sidera", font_size=54, color=TEXT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(4.6), SLIDE_W, Inches(0.6),
             "Your AI marketing team that works in Slack.",
             font_size=22, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(6.2), SLIDE_W, Inches(0.4),
             "Confidential \u2014 February 2026",
             font_size=12, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 1, TOTAL_SLIDES)


# ── SLIDE 2: The Problem ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0), Inches(0.8), SLIDE_W, Inches(0.8),
             "The Problem", font_size=44, color=TEXT, bold=True,
             alignment=PP_ALIGN.CENTER)

cards = [
    ("Data Overload", "Google Ads, Meta, BigQuery, GA4 \u2014 your team is drowning in dashboards. Insights get missed.", RED),
    ("Misaligned Incentives", "Ad platforms optimize for their revenue, not yours. Platform-reported ROAS inflates reality by 20-40%.", ORANGE),
    ("Human Bottlenecks", "Your team can't monitor campaigns 24/7. Budget overruns, creative fatigue, and anomalies go unnoticed for days.", RED),
]
for i, (title, body, color) in enumerate(cards):
    left = Inches(1.0) + i * Inches(3.9)
    card_with_text(slide, left, Inches(2.2), Inches(3.5), Inches(3.5),
                   title, body, title_color=color)
add_slide_number(slide, 2, TOTAL_SLIDES)


# ── SLIDE 3: The Vision ───────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(2),
             "What if you could hire\nan AI marketing team\nthat works in Slack?",
             font_size=44, color=TEXT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(2),
             "Each member has a role, expertise, memory, and principles.\n"
             "They analyze your data, surface insights, and recommend actions.\n"
             "Nothing happens without your approval.",
             font_size=20, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 3, TOTAL_SLIDES)


# ── SLIDE 4: How It Works ─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0), Inches(0.8), SLIDE_W, Inches(0.8),
             "How It Works", font_size=44, color=TEXT, bold=True,
             alignment=PP_ALIGN.CENTER)

steps = [
    ("1. Connect", "Google Ads, Meta,\nBigQuery, Drive"),
    ("2. Analyze", "AI roles run 19 skills\ndaily, pulling live data"),
    ("3. Recommend", "Insights delivered to\nSlack with context"),
    ("4. You Decide", "Approve or Reject \u2014\nevery action is gated"),
]
for i, (title, body) in enumerate(steps):
    left = Inches(0.8) + i * Inches(3.15)
    bc = ACCENT if i == 3 else None
    card_with_text(slide, left, Inches(2.5), Inches(2.8), Inches(2.5),
                   title, body, border_color=bc)
    if i < 3:
        add_text_box(slide, left + Inches(2.85), Inches(3.4),
                     Inches(0.4), Inches(0.5),
                     "\u2192", font_size=28, color=ACCENT_LIGHT,
                     alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 4, TOTAL_SLIDES)


# ── SLIDE 5: Meet the Team ────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0), Inches(0.6), SLIDE_W, Inches(0.8),
             "Meet Your AI Team", font_size=44, color=TEXT, bold=True,
             alignment=PP_ALIGN.CENTER)

# Marketing dept label
add_text_box(slide, Inches(1), Inches(1.6), Inches(7), Inches(0.4),
             "MARKETING DEPARTMENT", font_size=12, color=ACCENT_LIGHT, bold=True,
             alignment=PP_ALIGN.CENTER)

# Manager box
card_with_text(slide, Inches(2.5), Inches(2.1), Inches(4), Inches(1.2),
               "Head of Marketing", "Portfolio strategy & executive briefings",
               border_color=ACCENT)

# Sub-roles
sub_roles = [
    ("Media Buyer", "11 skills"),
    ("Reporting Analyst", "2 skills"),
    ("Strategist", "2 skills"),
]
for i, (name, desc) in enumerate(sub_roles):
    left = Inches(1.0) + i * Inches(2.5)
    card_with_text(slide, left, Inches(3.9), Inches(2.2), Inches(1.2),
                   name, desc)

# IT dept
add_text_box(slide, Inches(8.5), Inches(1.6), Inches(4), Inches(0.4),
             "IT DEPARTMENT", font_size=12, color=BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)
card_with_text(slide, Inches(9), Inches(2.1), Inches(3), Inches(1.2),
               "Head of IT", "Self-monitoring & diagnostics",
               title_color=BLUE, border_color=BLUE)

add_text_box(slide, Inches(0), Inches(5.8), SLIDE_W, Inches(0.6),
             "The manager delegates to sub-roles, synthesizes their work, and delivers one unified briefing.",
             font_size=16, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 5, TOTAL_SLIDES)


# ── SLIDE 6: DEMO - Talk to team ──────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, RGBColor(0x0C, 0x12, 0x10))
add_text_box(slide, Inches(0), Inches(1.8), SLIDE_W, Inches(0.5),
             "LIVE DEMO", font_size=16, color=ACCENT_LIGHT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(2.5), SLIDE_W, Inches(1),
             "Talk to the Team", font_size=48, color=TEXT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(4.0), SLIDE_W, Inches(0.5),
             "@Project Sidera talk to the media buyer",
             font_size=22, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER,
             font_name="Consolas")
add_text_box(slide, Inches(0), Inches(5.0), SLIDE_W, Inches(1),
             "Ask about campaign performance.\nFollow up in the thread.",
             font_size=18, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 6, TOTAL_SLIDES)


# ── SLIDE 7: Approval Gate ────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0), Inches(0.8), SLIDE_W, Inches(0.8),
             "Nothing Happens Without You", font_size=44, color=TEXT, bold=True,
             alignment=PP_ALIGN.CENTER)

# Mock approval card
card = add_rounded_rect(slide, Inches(3.2), Inches(2.2), Inches(7), Inches(3.2),
                        border_color=RGBColor(0x30, 0x30, 0x40))
add_text_box(slide, Inches(3.5), Inches(2.4), Inches(6.4), Inches(0.4),
             "Performance Media Buyer recommends:", font_size=16, color=TEXT, bold=True)
add_text_box(slide, Inches(3.5), Inches(2.9), Inches(6.4), Inches(1.5),
             'Pause campaign "Brand - Exact Match"\n'
             "CPA $47.20, 2.3x above target. Backend ROAS 0.8x.\n"
             "Risk: Low \u2014 $120/day with declining performance over 7 days",
             font_size=15, color=TEXT_DIM)

# Approve button
approve = add_rounded_rect(slide, Inches(3.5), Inches(4.5), Inches(1.6), Inches(0.5),
                           fill_color=GREEN)
add_text_box(slide, Inches(3.5), Inches(4.52), Inches(1.6), Inches(0.45),
             "Approve", font_size=16, color=BLACK, bold=True,
             alignment=PP_ALIGN.CENTER)

# Reject button
reject = add_rounded_rect(slide, Inches(5.3), Inches(4.5), Inches(1.6), Inches(0.5),
                           fill_color=BG_CARD, border_color=RED)
add_text_box(slide, Inches(5.3), Inches(4.52), Inches(1.6), Inches(0.45),
             "Reject", font_size=16, color=RED, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.8),
             "Every write operation \u2014 budget changes, campaign pauses, bid adjustments \u2014\n"
             "requires explicit human approval before execution.",
             font_size=17, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 7, TOTAL_SLIDES)


# ── SLIDE 8: DEMO - Approval ─────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, RGBColor(0x0C, 0x12, 0x10))
add_text_box(slide, Inches(0), Inches(1.8), SLIDE_W, Inches(0.5),
             "LIVE DEMO", font_size=16, color=ACCENT_LIGHT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(2.5), SLIDE_W, Inches(1),
             "Approve a Recommendation", font_size=48, color=TEXT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(4.2), SLIDE_W, Inches(1.5),
             "Agent recommends a change.\nClick Approve in Slack.\nWatch it execute in real-time.",
             font_size=20, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 8, TOTAL_SLIDES)


# ── SLIDE 9: It Learns ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0), Inches(0.8), SLIDE_W, Inches(0.8),
             "It Gets Smarter Every Day", font_size=44, color=TEXT, bold=True,
             alignment=PP_ALIGN.CENTER)

learn_cards = [
    ("Decisions", "Remembers every approval and rejection. Learns what you actually want over time.", ACCENT_LIGHT),
    ("Patterns", "Detects recurring anomalies, budget trends, and seasonal effects across runs.", ACCENT_LIGHT),
    ("Lessons", '"I tried X, it failed because Y." After each run, the agent reflects on what it could do better.', ACCENT_LIGHT),
]
for i, (title, body, color) in enumerate(learn_cards):
    left = Inches(1.0) + i * Inches(3.9)
    card_with_text(slide, left, Inches(2.2), Inches(3.5), Inches(3.0),
                   title, body, title_color=color)

add_text_box(slide, Inches(0), Inches(5.8), SLIDE_W, Inches(0.5),
             "Memory persists across sessions. Hot memories auto-inject into context. Cold memories are searchable.",
             font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 9, TOTAL_SLIDES)


# ── SLIDE 10: Morning Briefing ────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0), Inches(0.8), SLIDE_W, Inches(0.8),
             "Your Morning Briefing", font_size=44, color=TEXT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(1.8), SLIDE_W, Inches(0.5),
             "Every weekday at 9 AM, the Head of Marketing:",
             font_size=20, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

briefing_steps = [
    ("1", "Runs its own\nexecutive summary"),
    ("2", "Decides which\nsub-roles to activate"),
    ("3", "Collects their\nanalyses"),
    ("4", "Synthesizes one\nunified briefing"),
]
for i, (num, desc) in enumerate(briefing_steps):
    left = Inches(1.2) + i * Inches(2.9)
    add_rounded_rect(slide, left, Inches(2.8), Inches(2.5), Inches(2.2))
    add_text_box(slide, left, Inches(3.0), Inches(2.5), Inches(0.6),
                 num, font_size=36, color=BLUE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), Inches(3.7), Inches(2.2), Inches(1),
                 desc, font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(5.5), SLIDE_W, Inches(0.6),
             "Cross-channel insights that individual roles would miss.\nConflicts flagged. Actions prioritized.",
             font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 10, TOTAL_SLIDES)


# ── SLIDE 11: DEMO - Run HoM ─────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, RGBColor(0x0C, 0x12, 0x10))
add_text_box(slide, Inches(0), Inches(1.8), SLIDE_W, Inches(0.5),
             "LIVE DEMO", font_size=16, color=ACCENT_LIGHT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(2.5), SLIDE_W, Inches(1),
             "The Morning Briefing", font_size=48, color=TEXT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(4.0), SLIDE_W, Inches(0.5),
             "/sidera run role:head_of_marketing",
             font_size=22, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER,
             font_name="Consolas")
add_text_box(slide, Inches(0), Inches(5.0), SLIDE_W, Inches(0.5),
             "Watch the manager delegate, collect, and synthesize.",
             font_size=18, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 11, TOTAL_SLIDES)


# ── SLIDE 12: Source of Truth ─────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0), Inches(0.8), SLIDE_W, Inches(0.8),
             "Your Data. Your Truth.", font_size=44, color=TEXT, bold=True,
             alignment=PP_ALIGN.CENTER)

# BigQuery card
card_with_text(slide, Inches(1.5), Inches(2.3), Inches(4.5), Inches(2.8),
               "BigQuery Backend",
               "Your own attribution model.\nReal revenue, real orders, real ROAS.\n\n"
               "\u2705 SOURCE OF TRUTH",
               title_color=GREEN, border_color=ACCENT)

add_text_box(slide, Inches(6.1), Inches(3.2), Inches(1), Inches(0.5),
             "vs.", font_size=24, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# Platform card
card_with_text(slide, Inches(7.3), Inches(2.3), Inches(4.5), Inches(2.8),
               "Platform-Reported",
               "Google and Meta self-report.\nTypically inflated 20-40%\nvs. backend attribution.\n\n"
               "\u26a0\ufe0f CROSS-REFERENCED",
               title_color=ORANGE)

add_text_box(slide, Inches(0), Inches(5.7), SLIDE_W, Inches(0.6),
             "The agent calculates the inflation ratio for every campaign.\nDecisions based on reality, not platform spin.",
             font_size=15, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 12, TOTAL_SLIDES)


# ── SLIDE 13: It Scales ───────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0), Inches(0.8), SLIDE_W, Inches(0.8),
             "It Scales to Any Domain", font_size=44, color=TEXT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(1.8), SLIDE_W, Inches(0.6),
             "The core framework is domain-agnostic. Marketing is the first department.",
             font_size=20, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

domains = [
    ("Marketing", "Live now", ACCENT_LIGHT, ACCENT),
    ("Finance", "Budget & forecasting", TEXT, None),
    ("Sales", "Pipeline & leads", TEXT, None),
    ("Operations", "Supply chain", TEXT, None),
]
for i, (name, desc, color, border) in enumerate(domains):
    left = Inches(1.5) + i * Inches(2.8)
    card_with_text(slide, left, Inches(3.0), Inches(2.4), Inches(1.8),
                   name, desc, title_color=color, border_color=border)

add_text_box(slide, Inches(0), Inches(5.5), SLIDE_W, Inches(0.5),
             "Same pattern: connect APIs, teach skills via YAML, approve in Slack, log everything.",
             font_size=14, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 13, TOTAL_SLIDES)


# ── SLIDE 14: What's Next ─────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0), Inches(0.8), SLIDE_W, Inches(0.8),
             "What's Next", font_size=44, color=TEXT, bold=True,
             alignment=PP_ALIGN.CENTER)

items = [
    ("Meta Ads integration", "Connector built, E2E testing underway"),
    ("Graduated trust", "Auto-approve low-risk actions matching pre-set rules"),
    ("100+ skill library", "Agents teach themselves new skills (with your approval)"),
    ("Voice meetings", "Department heads join calls as voice participants"),
    ("Production deployment", "Railway, always-on, morning briefings on autopilot"),
]
for i, (title, desc) in enumerate(items):
    y = Inches(2.2) + i * Inches(0.9)
    add_text_box(slide, Inches(2), y, Inches(3), Inches(0.4),
                 title, font_size=20, color=TEXT, bold=True)
    add_text_box(slide, Inches(5.2), y, Inches(6), Inches(0.4),
                 "\u2014  " + desc, font_size=16, color=TEXT_DIM)
add_slide_number(slide, 14, TOTAL_SLIDES)


# ── SLIDE 15: Thank You ───────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(
        LOGO_PATH,
        SLIDE_W / 2 - Inches(0.9), Inches(1.5),
        width=Inches(1.8), height=Inches(1.8)
    )
add_text_box(slide, Inches(0), Inches(3.6), SLIDE_W, Inches(1),
             "Thank You", font_size=54, color=TEXT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(4.6), SLIDE_W, Inches(0.5),
             "Questions?", font_size=22, color=TEXT_DIM,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(6.2), SLIDE_W, Inches(0.4),
             "\u25bc  Technical deep-dive follows  \u25bc",
             font_size=12, color=RGBColor(0x40, 0x40, 0x50),
             alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 15, TOTAL_SLIDES)


# ── SLIDE 16: APPENDIX DIVIDER ────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, RGBColor(0x0D, 0x0B, 0x14))
if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(
        LOGO_PATH,
        SLIDE_W / 2 - Inches(0.6), Inches(1.5),
        width=Inches(1.2), height=Inches(1.2)
    )
add_text_box(slide, Inches(0), Inches(3.0), SLIDE_W, Inches(0.4),
             "APPENDIX", font_size=14, color=PURPLE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(3.5), SLIDE_W, Inches(1),
             "Technical Deep Dive", font_size=48, color=TEXT, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0), Inches(4.6), SLIDE_W, Inches(0.5),
             "Architecture, systems, and the decisions behind them.",
             font_size=18, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 16, TOTAL_SLIDES)


# ── SLIDE 17: System Architecture ─────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.7),
             "System Architecture", font_size=40, color=TEXT, bold=True)

# Core stack card
card_with_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.5),
               "Core Stack",
               "Agent: Anthropic Claude API (Python SDK)\n"
               "Orchestration: Inngest durable functions\n"
               "Database: PostgreSQL (Supabase)\n"
               "Cache: Redis (Upstash)\n"
               "API: FastAPI\n"
               "Notifications: Slack Bolt SDK\n"
               "Deploy: Railway")

# Numbers card
add_rounded_rect(slide, Inches(6.8), Inches(1.6), Inches(5.7), Inches(4.5))
add_text_box(slide, Inches(7.1), Inches(1.8), Inches(5), Inches(0.4),
             "By the Numbers", font_size=20, color=ACCENT_LIGHT, bold=True)

numbers = [
    ("8", "Connectors"), ("55", "MCP Tools"),
    ("19", "Skills"), ("15", "Workflows"),
    ("98", "DB Methods"), ("2845+", "Tests"),
]
for i, (num, label) in enumerate(numbers):
    col = i % 3
    row = i // 3
    x = Inches(7.2) + col * Inches(1.8)
    y = Inches(2.5) + row * Inches(1.6)
    add_text_box(slide, x, y, Inches(1.6), Inches(0.6),
                 num, font_size=32, color=ACCENT_LIGHT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(0.55), Inches(1.6), Inches(0.4),
                 label, font_size=13, color=TEXT_DIM,
                 alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 17, TOTAL_SLIDES)


# ── SLIDE 18: Model Routing ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0.8), Inches(0.6), Inches(10), Inches(0.7),
             "Intelligent Model Routing", font_size=40, color=TEXT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Each briefing uses three AI models. Total: ~$0.52 vs $1.50-3.00 single-model.",
             font_size=16, color=TEXT_DIM)

phases = [
    ("Phase 1: Haiku", "Data collection. Pulls raw numbers from Google Ads, Meta, BigQuery via tool calls.", "~$0.02", GREEN),
    ("Phase 2: Sonnet", "Tactical analysis. Identifies anomalies, trends, and actionable recommendations.", "~$0.15", BLUE),
    ("Phase 3: Opus", "Strategic layer. Cross-channel insights, portfolio thinking. Skipped on stable days.", "~$0.35", PURPLE),
]
for i, (title, desc, cost, color) in enumerate(phases):
    left = Inches(0.8) + i * Inches(4.1)
    card_with_text(slide, left, Inches(2.2), Inches(3.7), Inches(3.5),
                   title, desc + "\n\n" + cost, title_color=color)

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(10), Inches(0.4),
             "Opus auto-skips when volatility < 10%. Haiku failure falls back to cached data.",
             font_size=13, color=TEXT_DIM)
add_slide_number(slide, 18, TOTAL_SLIDES)


# ── SLIDE 19: Skill System ────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.7),
             "The Skill System", font_size=40, color=TEXT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Skills are YAML files on disk \u2014 version controlled, human-readable, composable.",
             font_size=16, color=TEXT_DIM)

# Hierarchy
card_with_text(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.2),
               "Three-Level Hierarchy",
               "Department \u2192 shared context\n"
               '(e.g. "trust backend over platform data")\n\n'
               "Role \u2192 persona, principles, connectors, schedule\n\n"
               "Skill \u2192 specific task\n"
               "(anomaly detection, budget pacing, etc.)\n\n"
               "Context flows down: dept context + role persona\n"
               "+ skill supplement compose into one prompt.")

# Skills grid
add_rounded_rect(slide, Inches(6.8), Inches(2.0), Inches(5.7), Inches(4.2))
add_text_box(slide, Inches(7.1), Inches(2.2), Inches(5), Inches(0.4),
             "Media Buyer: 11 Skills", font_size=18, color=ACCENT_LIGHT, bold=True)

skills = [
    "Anomaly Detector", "Creative Analysis", "Budget Reallocation",
    "Budget Pacing", "Creative Fatigue", "Bid Strategy Review",
    "Search Term Audit", "Dayparting Analysis", "Geo Performance",
    "Landing Page Analysis", "Platform Health Check"
]
for i, skill in enumerate(skills):
    col = i % 2
    row = i // 2
    x = Inches(7.1) + col * Inches(2.7)
    y = Inches(2.8) + row * Inches(0.5)
    add_text_box(slide, x, y, Inches(2.5), Inches(0.4),
                 "\u2022  " + skill, font_size=12, color=TEXT_DIM)
add_slide_number(slide, 19, TOTAL_SLIDES)


# ── SLIDE 20: Manager Delegation ──────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0.8), Inches(0.6), Inches(10), Inches(0.7),
             "Manager Delegation Pipeline", font_size=40, color=TEXT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             'A manager is just a role with a "manages" field. No new tables, no new abstractions.',
             font_size=16, color=TEXT_DIM)

mgr_steps = [
    ("Own Skills", "Runs its own\nbriefing skills first"),
    ("Delegation", "LLM decides which\nsub-roles to activate"),
    ("Sub-Role Exec", "Each sub-role runs\nwith full tool access"),
    ("Synthesis", "Combines all outputs\ninto unified briefing"),
]
for i, (title, desc) in enumerate(mgr_steps):
    left = Inches(0.8) + i * Inches(3.15)
    card_with_text(slide, left, Inches(2.2), Inches(2.8), Inches(2.2),
                   title, desc, title_color=BLUE)
    if i < 3:
        add_text_box(slide, left + Inches(2.85), Inches(3.0),
                     Inches(0.4), Inches(0.5),
                     "\u2192", font_size=28, color=BLUE,
                     alignment=PP_ALIGN.CENTER)

card_with_text(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(1.5),
               "Delegation Fallback",
               "If the delegation LLM call fails, activates ALL sub-roles. Safe default \u2014 never loses data.")
card_with_text(slide, Inches(6.8), Inches(5.0), Inches(5.5), Inches(1.5),
               "Recursive Managers",
               "Managers can manage other managers. Depth limit of 3 prevents runaway chains.")
add_slide_number(slide, 20, TOTAL_SLIDES)


# ── SLIDE 21: Safety ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.7),
             "Safety Architecture", font_size=40, color=TEXT, bold=True)

card_with_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(3.5),
               "Write Operation Safety",
               "\u2022 50% budget cap \u2014 no change exceeds 50% of current\n"
               "\u2022 Double-execution prevention via executed_at field\n"
               "\u2022 Previous values stored for rollback\n"
               "\u2022 Every write needs explicit human approval\n"
               "\u2022 Full audit trail \u2014 who, when, what changed")

card_with_text(slide, Inches(6.8), Inches(1.6), Inches(5.7), Inches(3.5),
               "Graduated Trust (Auto-Execute)",
               "\u2022 YAML rules per role define what can auto-execute\n"
               "\u2022 10 condition operators (less_than, between, etc.)\n"
               "\u2022 Constraints: daily limits, cooldowns, platform\n"
               "\u2022 Global kill switch (default: OFF)\n"
               "\u2022 Skill proposals NEVER auto-execute (hard block)")

card_with_text(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.3),
               "Information Security",
               "Four clearance levels: PUBLIC \u2192 INTERNAL \u2192 CONFIDENTIAL \u2192 RESTRICTED. "
               "Skills gated by min clearance. Roles have clearance levels. Inter-agent data filtered.",
               title_color=ORANGE)
add_slide_number(slide, 21, TOTAL_SLIDES)


# ── SLIDE 22: Conversational Mode ─────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.7),
             "Conversational Mode", font_size=40, color=TEXT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Every role is dual-mode: autonomous (scheduled) and conversational (Slack threads).",
             font_size=16, color=TEXT_DIM)

card_with_text(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3.2),
               "How It Works",
               "\u2022 @mention or /sidera chat starts a thread\n"
               "\u2022 RoleRouter identifies role (regex + Haiku)\n"
               "\u2022 Thread pinned to one role via DB\n"
               "\u2022 Each reply = new agent turn with tools\n"
               "\u2022 Stateless: history from Slack API each turn")

card_with_text(slide, Inches(6.8), Inches(2.0), Inches(5.7), Inches(3.2),
               "Write Ops in Conversations",
               "\u2022 Agent generates JSON recommendations\n"
               "\u2022 System posts Approve/Reject in-thread\n"
               "\u2022 On approval, executes via connector\n"
               "\u2022 Same safety pipeline as briefings")

# Limits
limits = [("20", "Max turns"), ("24h", "Timeout"), ("$5", "Cost cap")]
for i, (num, label) in enumerate(limits):
    left = Inches(1.5) + i * Inches(3.8)
    add_rounded_rect(slide, left, Inches(5.6), Inches(2.5), Inches(1.2))
    add_text_box(slide, left, Inches(5.7), Inches(2.5), Inches(0.5),
                 num, font_size=28, color=ACCENT_LIGHT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(6.2), Inches(2.5), Inches(0.4),
                 label, font_size=13, color=TEXT_DIM,
                 alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 22, TOTAL_SLIDES)


# ── SLIDE 23: Memory System ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.7),
             "Memory System", font_size=40, color=TEXT, bold=True)

card_with_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(3.0),
               "Tiered Architecture",
               "\u2022 Hot memories (\u226490 days) \u2014 auto-injected, 2000 token cap\n"
               "\u2022 Cold archive (>90 days) \u2014 never deleted, searchable\n"
               "\u2022 Weekly consolidation (Haiku, Sunday 4 AM)\n"
               "\u2022 Memory versioning via supersedes_id chains")

card_with_text(slide, Inches(6.8), Inches(1.6), Inches(5.7), Inches(3.0),
               "Memory Types",
               "Decision \u2022 Anomaly \u2022 Pattern \u2022 Insight \u2022 Lesson \u2022 Relationship\n\n"
               "Decisions from approvals. Anomalies from keywords.\n"
               "Lessons from post-run reflection. No LLM needed.")

card_with_text(slide, Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.5),
               "Post-Run Reflection",
               "After every role run, a Haiku call (~$0.01) asks: \"What was hard? What data was missing? "
               "What would you do differently?\" Answers saved as lesson memories. Compounding value.",
               title_color=GREEN)
add_slide_number(slide, 23, TOTAL_SLIDES)


# ── SLIDE 24: Skill & Role Evolution ──────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.7),
             "Self-Improving Agents", font_size=40, color=TEXT, bold=True)

card_with_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(3.2),
               "Skill Evolution",
               "Agents propose changes to their own skills.\n"
               "Same pipeline as budget changes:\n\n"
               "Agent proposes \u2192 Approval queue \u2192\n"
               "Human reviews diff in Slack \u2192\n"
               "Approve/Reject \u2192 Written to DB\n\n"
               "FORBIDDEN: requires_approval, manages, is_active")

card_with_text(slide, Inches(6.8), Inches(1.6), Inches(5.7), Inches(3.2),
               "Role Evolution",
               "Department heads propose new roles.\n\n"
               "\u2022 Only managers can propose\n"
               "\u2022 Department-scoped\n"
               "\u2022 New roles auto-added to manages list\n"
               "\u2022 NEVER auto-executes")

card_with_text(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.5),
               "Dynamic Org Chart",
               "YAML on disk = seed data. Database = override layer. /sidera org Slack commands + REST API. "
               "All mutations audit-logged. Soft delete. Graceful DB fallback.")
add_slide_number(slide, 24, TOTAL_SLIDES)


# ── SLIDE 25: Connectors ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0.8), Inches(0.6), Inches(10), Inches(0.7),
             "Connectors & Production Hardening", font_size=40, color=TEXT, bold=True)

connectors = [
    ("Google Ads", "7 read + 6 write"),
    ("Meta Ads", "7 read + 6 write"),
    ("BigQuery", "7 methods (source of truth)"),
    ("Google Drive", "13 methods (Docs/Sheets/Slides)"),
    ("Slack", "16 methods"),
    ("Recall.ai", "Meeting bot lifecycle"),
]

add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0))
add_text_box(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.4),
             "6 Data Connectors", font_size=20, color=ACCENT_LIGHT, bold=True)
for i, (name, desc) in enumerate(connectors):
    y = Inches(2.3) + i * Inches(0.5)
    add_text_box(slide, Inches(1.1), y, Inches(2.2), Inches(0.4),
                 name, font_size=14, color=TEXT, bold=True)
    add_text_box(slide, Inches(3.3), y, Inches(2.8), Inches(0.4),
                 desc, font_size=13, color=TEXT_DIM)

card_with_text(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.0),
               "Production Hardening",
               "\u2022 Retry with exponential backoff + jitter (all connectors)\n"
               "\u2022 Fernet token encryption (enc: prefix convention)\n"
               "\u2022 Dead letter queue (failed_runs table for replay)\n"
               "\u2022 Proactive token refresh (5 AM cron, 7-day lookahead)\n"
               "\u2022 Sentry error monitoring in all workflows\n"
               "\u2022 Rate limiting at API level\n"
               "\u2022 Redis caching (@cached decorator, 2h TTL)")
add_slide_number(slide, 25, TOTAL_SLIDES)


# ── SLIDE 26: Voice Meetings ──────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_watermark(slide)
add_text_box(slide, Inches(0.8), Inches(0.6), Inches(10), Inches(0.7),
             "Listen-Only Meeting Participation", font_size=40, color=TEXT, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
             "Manager roles can join live video calls as listen-only participants.",
             font_size=16, color=TEXT_DIM)

voice_steps = [
    ("Recall.ai", "Bot joins\nvia URL"),
    ("Webhooks", "Transcript\ncapture"),
    ("Claude", "Post-call\nanalysis"),
    ("Manager", "Delegation\npipeline"),
]
for i, (title, desc) in enumerate(voice_steps):
    left = Inches(0.8) + i * Inches(3.15)
    card_with_text(slide, left, Inches(2.2), Inches(2.8), Inches(1.8),
                   title, desc, title_color=BLUE)
    if i < 3:
        add_text_box(slide, left + Inches(2.85), Inches(2.8),
                     Inches(0.4), Inches(0.5),
                     "\u2192", font_size=28, color=BLUE,
                     alignment=PP_ALIGN.CENTER)

card_with_text(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.2),
               "In-Call",
               "~1.7s end-to-end latency\n"
               "500ms debounce\n"
               "Concurrency guards prevent overlapping\n"
               "Self-echo filtering\n"
               "200-char response cap")

card_with_text(slide, Inches(6.8), Inches(4.5), Inches(5.7), Inches(2.2),
               "Post-Call",
               "Transcript summary \u2192 action items \u2192\n"
               "manager delegation pipeline\n\n"
               "Automatic follow-up via Slack approval flow\n"
               "Limits: 2h max, $10/session cap")
add_slide_number(slide, 26, TOTAL_SLIDES)


# ── SLIDE 27: The Big Picture ─────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(
        LOGO_PATH,
        SLIDE_W / 2 - Inches(0.6), Inches(0.8),
        width=Inches(1.2), height=Inches(1.2)
    )
add_text_box(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(2.5),
             "Sidera is not a chatbot.\n\n"
             "It's a framework for building AI workforces \u2014\n"
             "teams of specialized agents with roles, memory,\n"
             "principles, and safety controls, orchestrated through\n"
             "durable workflows and gated by human approval.",
             font_size=24, color=TEXT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.8),
             "The accumulated expertise in each skill \u2014 the context, examples,\n"
             "guidelines, and learned memories \u2014 is the product.\n"
             "The framework just runs it.",
             font_size=16, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0), Inches(6.5), SLIDE_W, Inches(0.4),
             "2 departments  \u2022  5 roles  \u2022  19 skills  \u2022  55 tools  \u2022  15 workflows  \u2022  8 connectors  \u2022  2845+ tests",
             font_size=11, color=RGBColor(0x40, 0x40, 0x50),
             alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 27, TOTAL_SLIDES)


# ── Save ──────────────────────────────────────────────────────
output_path = os.path.join(SCRIPT_DIR, "Project Sidera - Pitch Deck.pptx")
prs.save(output_path)
print(f"Saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
